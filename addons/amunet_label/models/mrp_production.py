# -*- coding: utf-8 -*-
# Genera las etiquetas de caja en PPTX (tabloide 11x17) directamente desde
# la orden de fabricacion, usando python-pptx (instalado en la imagen) y las
# plantillas PLANTILLA_S/H/P/M.pptx que viven en static/templates/ del modulo.
# Los datos (producto, REF, lote, caducidad, contenedor, accesorio, registro
# sanitario, presentaciones autorizadas) se leen del producto y de la MO en
# Odoo; NO depende de ningun Excel externo.
import io
import os
import re
import base64

from odoo import models, fields, api, _
from odoo.exceptions import UserError

# Namespace DrawingML: los textos de las formas viven en elementos <a:t>.
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class MrpProduction(models.Model):
    _inherit = 'mrp.production'

    def _etiqueta_datos(self):
        """Datos fijos de la etiqueta (sin la cantidad por caja).
        Devuelve (subtipo, lot_name, placeholders)."""
        self.ensure_one()
        prod = self.product_id
        tmpl = prod.product_tmpl_id
        # El subtipo puede venir vacio: solo hace falta si la presentacion usa
        # caja generica (etiqueta grande). Para cajas pre-impresas (Plantilla A)
        # no se necesita. La validacion se hace al construir, no aqui.
        subtipo = tmpl.etiqueta_subtipo

        # Lote producido: en Odoo 19 la MO no tiene lot_producing_id; el lote
        # sale de las lineas de movimiento del producto terminado.
        lote = self.move_finished_ids.move_line_ids.filtered(
            lambda ml: ml.lot_id and ml.product_id == prod
        ).mapped('lot_id')[:1]
        lot_name = lote.name if lote else self.name
        # Caducidad = la de la ORDEN (amunet_expiration_text, ej "2028-05"), que es
        # la real y la que se ve en el plan. NO la expiration_date del lote, que a
        # veces trae la fecha de creacion del lote (bug observado en Calprotectina:
        # el lote decia 2026-07 y la orden 2028-05).
        caducidad = getattr(self, 'amunet_expiration_text', '') or ''
        if not caducidad:
            exp = lote.expiration_date if lote else False
            caducidad = str(exp)[:7] if exp else '????-??'

        # Contenedor y accesorio van justo despues de "{{N}} " en la etiqueta
        # (ej. "-20 tubos..."). La primera letra debe ir en MAYUSCULA aunque en
        # la ficha del producto se haya capturado en minuscula.
        def _cap(s):
            s = s or ''
            return s[0].upper() + s[1:] if s else s

        datos = {
            '{{PRODUCTO}}': tmpl.nombre_etiqueta or prod.name,
            '{{REF}}': prod.default_code or '',
            '{{LOT}}': lot_name,
            '{{LOTE}}': lot_name,  # Plantilla A usa {{LOTE}} en vez de {{LOT}}
            '{{CADUCIDAD}}': caducidad,
            '{{CONTENEDOR}}': _cap(tmpl.etiqueta_contenedor),
            '{{ACCESORIO}}': _cap(tmpl.etiqueta_accesorio),
            '{{REG_SANITARIO}}': tmpl.registro_sanitario or '',
        }
        return subtipo, lot_name, datos

    def _etiqueta_clonar_slide(self, prs, src):
        """Duplica la diapositiva `src` dentro del MISMO pptx (mismo paquete,
        asi las imagenes/diseno se conservan) y devuelve la nueva."""
        import copy
        from pptx.oxml.ns import qn
        dst = prs.slides.add_slide(src.slide_layout)
        # Quitar los placeholders que el layout agrega por defecto.
        for shp in list(dst.shapes):
            shp._element.getparent().remove(shp._element)
        # Copiar cada forma de la diapositiva origen.
        for shp in src.shapes:
            dst.shapes._spTree.append(copy.deepcopy(shp._element))
        # Reapuntar relaciones de imagenes/medios (r:embed / r:link) a las
        # mismas partes ya presentes en el paquete.
        for attr in (qn('r:embed'), qn('r:link')):
            for el in dst.shapes._spTree.iter():
                rid = el.get(attr)
                if rid and rid in src.part.rels:
                    rel = src.part.rels[rid]
                    nuevo = dst.part.relate_to(rel.target_part, rel.reltype)
                    el.set(attr, nuevo)
        return dst

    # Rejilla ETIQUETA GRANDE (S/H/P/M) 3.53x2.72": 3 columnas x 6 filas = 18/hoja.
    _GRID_COLS_IN = (0.19, 3.73, 7.27)
    _GRID_ROW0_IN = 0.25
    _GRID_ROW_PITCH_IN = 2.74
    _GRID_ROWS = 6

    # Rejilla PLANTILLA A (etiqueta chica 1.05x0.63" para cajas pre-impresas):
    # 10 columnas x 25 filas = 250/hoja tabloide.
    _GRID_A_COLS = tuple(round(0.19 + 1.06 * i, 3) for i in range(10))
    _GRID_A_ROW0 = 0.25
    _GRID_A_PITCH = 0.65
    _GRID_A_ROWS = 25

    def _etiqueta_importar_medias(self, prs_dst, src_part, grupo_el):
        """Copia al paquete destino las partes de media (imagen/svg) que usa un
        grupo importado de OTRA plantilla (ej. el SVG de la Plantilla A al meterlo
        en el archivo de las etiquetas grandes). Devuelve {rid_original: (parte, reltype)}."""
        from pptx.oxml.ns import qn
        from pptx.opc.packuri import PackURI
        from pptx.opc.package import Part
        pkg = prs_dst.part.package
        existentes = {p.partname for p in pkg.iter_parts()}
        media = {}
        for attr in (qn('r:embed'), qn('r:link')):
            for node in grupo_el.iter():
                rid = node.get(attr)
                if not rid or rid in media or rid not in src_part.rels:
                    continue
                rel = src_part.rels[rid]
                src = rel.target_part
                ext = src.partname.ext
                i = 1
                while PackURI('/ppt/media/imgA%d.%s' % (i, ext)) in existentes:
                    i += 1
                pn = PackURI('/ppt/media/imgA%d.%s' % (i, ext))
                existentes.add(pn)
                media[rid] = (Part(pn, src.content_type, pkg, src.blob), rel.reltype)
        return media

    def _etiqueta_remap_medias(self, el, media, dst_part):
        """Reapunta las relaciones de un grupo a las medias ya importadas."""
        from pptx.oxml.ns import qn
        for attr in (qn('r:embed'), qn('r:link')):
            for node in el.iter():
                rid = node.get(attr)
                if rid in media:
                    parte, reltype = media[rid]
                    node.set(attr, dst_part.relate_to(parte, reltype))

    def _etiqueta_llenar_element(self, el, placeholders):
        """Reemplaza los {{...}} en todos los textos dentro de un elemento XML."""
        for t in el.iter('{%s}t' % _A_NS):
            if not t.text:
                continue
            nuevo = t.text
            for k, v in placeholders.items():
                nuevo = nuevo.replace(k, v)
            t.text = nuevo

    def _etiqueta_mover_grupo(self, grp_el, x_emu, y_emu, ref_el=None):
        """Mueve un grupo (p:grpSp) fijando su offset absoluto en la hoja.
        Blindaje: si el grupo copiado no trae su transform (a:xfrm) o su offset
        (a:off) — caso raro observado en generaciones grandes — se reconstruye
        desde `ref_el` (el grupo original de la plantilla, que siempre lo trae)
        para no romper la generacion de la etiqueta."""
        import copy as _copy
        from pptx.oxml.ns import qn
        xfrm = grp_el.find('%s/%s' % (qn('p:grpSpPr'), qn('a:xfrm')))
        if (xfrm is None or xfrm.find(qn('a:off')) is None) and ref_el is not None:
            ref_xfrm = ref_el.find('%s/%s' % (qn('p:grpSpPr'), qn('a:xfrm')))
            grp_sppr = grp_el.find(qn('p:grpSpPr'))
            if ref_xfrm is not None and grp_sppr is not None:
                if xfrm is not None:
                    grp_sppr.remove(xfrm)
                xfrm = _copy.deepcopy(ref_xfrm)
                grp_sppr.insert(0, xfrm)
        if xfrm is None:
            return
        off = xfrm.find(qn('a:off'))
        if off is None:
            return
        off.set('x', str(int(x_emu)))
        off.set('y', str(int(y_emu)))

    def _etiqueta_remap_rels(self, el, src_part, dst_part):
        """Reapunta las relaciones de imagen/medios (r:embed / r:link) de un
        elemento a la hoja destino, creando la relacion hacia la MISMA parte de
        imagen. Sin esto, las imagenes de la etiqueta no cargan en las hojas
        clonadas (2+) y el formato se ve roto."""
        from pptx.oxml.ns import qn
        for attr in (qn('r:embed'), qn('r:link')):
            for node in el.iter():
                rid = node.get(attr)
                if rid and rid in src_part.rels:
                    rel = src_part.rels[rid]
                    nuevo = dst_part.relate_to(rel.target_part, rel.reltype)
                    node.set(attr, nuevo)

    def _etiqueta_tile(self, prs, base, base_part, unidad_xml, posiciones,
                       hojas, hoja0, valores, media_a=None):
        """Coloca `valores` (una etiqueta c/u) copiando `unidad_xml` en las
        `hojas` a partir del indice `hoja0`, segun `posiciones` (rejilla).
        Si `media_a` viene, remapea medias importadas (Plantilla A); si no,
        remapea las relaciones ya presentes en `base_part` (etiqueta grande)."""
        import copy
        por_hoja = len(posiciones)
        for idx, vals in enumerate(valores):
            hoja = hojas[hoja0 + idx // por_hoja]
            x_emu, y_emu = posiciones[idx % por_hoja]
            nuevo = copy.deepcopy(unidad_xml)
            self._etiqueta_mover_grupo(nuevo, x_emu, y_emu, ref_el=unidad_xml)
            self._etiqueta_llenar_element(nuevo, vals)
            hoja.shapes._spTree.append(nuevo)
            if media_a is not None:
                self._etiqueta_remap_medias(nuevo, media_a, hoja.part)
            else:
                self._etiqueta_remap_rels(nuevo, base_part, hoja.part)

    def _etiqueta_construir_pptx(self, subtipo, datos, bloques):
        """Arma UN pptx tabloide con TODAS las etiquetas del plan. Cada bloque
        (una linea del plan) trae {'tipo': 'grande'|'A', 'n', 'cajas'}:
          - 'grande' -> etiqueta S/H/P/M (rejilla 3x6=18/hoja)
          - 'A'      -> Plantilla A chica para caja pre-impresa (rejilla 10x25=250/hoja)
        Todo en un solo archivo, reusando el mismo marco tabloide."""
        import copy
        import math
        from pptx import Presentation
        from pptx.util import Inches
        tdir = os.path.join(os.path.dirname(__file__), '..', 'static', 'templates')

        grandes = [(b['n'], b['cajas']) for b in bloques if b.get('tipo') != 'A']
        aes = [(b['n'], b['cajas']) for b in bloques if b.get('tipo') == 'A']
        if not grandes and not aes:
            grandes = [(0, 1)]

        # Base = plantilla grande (aporta marco + master + tema). Si no hay lineas
        # grandes usamos 'S' solo por su marco (identico en todas).
        if grandes and not subtipo:
            raise UserError(_(
                'El producto no tiene "Subtipo de etiqueta" y tiene presentaciones '
                'con caja general (etiqueta grande). Captura el subtipo en la ficha.'))
        sub_base = subtipo or 'S'
        ruta_g = os.path.join(tdir, 'PLANTILLA_%s.pptx' % sub_base)
        if not os.path.exists(ruta_g):
            raise UserError(_('No existe la plantilla PLANTILLA_%s.pptx en el modulo.') % sub_base)
        prs = Presentation(ruta_g)
        base = prs.slides[0]
        grupo_g = next((sh for sh in base.shapes if sh.shape_type == 6), None)
        if grupo_g is None:
            raise UserError(_('La plantilla PLANTILLA_%s.pptx no tiene la etiqueta como grupo.') % sub_base)
        unidad_g = copy.deepcopy(grupo_g._element)
        grupo_g._element.getparent().remove(grupo_g._element)
        base_part = base.part  # conserva las medias de la etiqueta grande

        # Grupo de la Plantilla A + import de su media (SVG) al paquete de salida.
        unidad_a = None
        media_a = None
        if aes:
            prs_a = Presentation(os.path.join(tdir, 'PLANTILLA_A.pptx'))
            grupo_a = next((sh for sh in prs_a.slides[0].shapes if sh.shape_type == 6), None)
            if grupo_a is None:
                raise UserError(_('PLANTILLA_A.pptx no tiene la etiqueta como grupo.'))
            unidad_a = copy.deepcopy(grupo_a._element)
            media_a = self._etiqueta_importar_medias(prs, prs_a.slides[0].part, unidad_a)

        # Valores por etiqueta.
        vals_g = []
        for n, cajas in grandes:
            for _i in range(cajas):
                vals_g.append(dict(datos, **{'{{N}}': str(n)}))
        vals_a = []
        if aes:
            a_base = {
                '{{REF}}': datos.get('{{REF}}', ''),
                '{{LOTE}}': datos.get('{{LOTE}}') or datos.get('{{LOT}}', ''),
                '{{CADUCIDAD}}': datos.get('{{CADUCIDAD}}', ''),
            }
            for n, cajas in aes:
                for _i in range(cajas):
                    vals_a.append(dict(a_base))

        # Rejillas (row-major).
        pos_g = [(Inches(c), Inches(self._GRID_ROW0_IN + r * self._GRID_ROW_PITCH_IN))
                 for r in range(self._GRID_ROWS) for c in self._GRID_COLS_IN]
        pos_a = [(Inches(c), Inches(self._GRID_A_ROW0 + r * self._GRID_A_PITCH))
                 for r in range(self._GRID_A_ROWS) for c in self._GRID_A_COLS]

        n_hojas_g = math.ceil(len(vals_g) / len(pos_g)) if vals_g else 0
        n_hojas_a = math.ceil(len(vals_a) / len(pos_a)) if vals_a else 0
        total_hojas = max(1, n_hojas_g + n_hojas_a)
        hojas = [base]
        for _h in range(1, total_hojas):
            hojas.append(self._etiqueta_clonar_slide(prs, base))  # marco solo

        if vals_g:
            self._etiqueta_tile(prs, base, base_part, unidad_g, pos_g, hojas, 0, vals_g)
        if vals_a:
            self._etiqueta_tile(prs, base, base_part, unidad_a, pos_a, hojas,
                                n_hojas_g, vals_a, media_a=media_a)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ==================================================================
    # ETIQUETAS DE BUFFER / SOLUCION DE CORRIMIENTO
    # Se generan junto a las de caja (mismo plan). El buffer es un move de
    # material de la orden con categoria "Semiterminado / Buffer" y plantilla
    # configurada. NAME = nombre del producto; LOT/CADUCIDAD = del lote del
    # buffer surtido (caducidad en año-mes).
    # ==================================================================
    def _etiqueta_buffers_de_orden(self, num_cajas=0):
        """Devuelve {plantilla: [valores,...]} con una entrada por vial a
        etiquetar. Un valor = dict de placeholders {{NAME}}/{{LOT}}/{{CADUCIDAD}}."""
        self.ensure_one()
        cat = self.env['product.category'].search(
            [('complete_name', '=', 'Semiterminado / Buffer')], limit=1)
        if not cat:
            return {}
        nombre = (self.product_id.product_tmpl_id.nombre_etiqueta or '').strip() \
            or self.product_id.name

        def _vals(lot, lot_name=''):
            cad = lot.expiration_date.strftime('%Y-%m') \
                if (lot and lot.expiration_date) else ''
            return {
                '{{NAME}}': nombre,
                '{{LOT}}': (lot.name if lot else lot_name) or '',
                '{{CADUCIDAD}}': cad,
            }

        specs = {}
        moves = self.move_raw_ids.filtered(
            lambda m: m.state != 'cancel'
            and m.product_id.categ_id == cat
            and m.product_id.product_tmpl_id.etiqueta_buffer_plantilla)
        for m in moves:
            tmpl = m.product_id.product_tmpl_id
            plantilla = tmpl.etiqueta_buffer_plantilla
            modo = tmpl.etiqueta_buffer_modo or 'por_vial'
            mls = m.move_line_ids.filtered(lambda l: (l.quantity or 0) > 0)
            if not mls:
                continue
            specs.setdefault(plantilla, [])
            if modo == 'por_caja':
                # 1 etiqueta por caja del plan (aunque la caja lleve varios
                # viales, ej. combo). Usa el lote surtido (primer move_line).
                vals = _vals(mls[0].lot_id, mls[0].lot_name or '')
                n = int(num_cajas) or int(round(sum(mls.mapped('quantity'))))
                specs[plantilla].extend([dict(vals) for _ in range(n)])
            else:  # por_vial: 1 por cada vial surtido
                for ml in mls:
                    qty = int(round(ml.quantity or 0))
                    vals = _vals(ml.lot_id, ml.lot_name or '')
                    specs[plantilla].extend([dict(vals) for _ in range(qty)])
        return specs

    def _etiqueta_construir_buffer_pptx(self, plantilla, valores):
        """Arma un PPTX de etiquetas de buffer con UNA plantilla, mosaico segun
        el tamaño de su etiqueta sobre la hoja 11x17."""
        import copy
        import math
        from pptx import Presentation
        tdir = os.path.join(os.path.dirname(__file__), '..', 'static', 'templates')
        ruta = os.path.join(tdir, 'PLANTILLA_BUFFER_%s.pptx' % plantilla)
        if not os.path.exists(ruta):
            raise UserError(_(
                'No existe la plantilla PLANTILLA_BUFFER_%s.pptx en el modulo.')
                % plantilla)
        prs = Presentation(ruta)
        base = prs.slides[0]
        grupo = next((sh for sh in base.shapes if sh.shape_type == 6), None)
        if grupo is None:
            raise UserError(_(
                'La plantilla PLANTILLA_BUFFER_%s.pptx no tiene la etiqueta '
                'como grupo.') % plantilla)
        lw, lh = int(grupo.width), int(grupo.height)
        x0, y0 = int(grupo.left), int(grupo.top)
        unidad = copy.deepcopy(grupo._element)
        grupo._element.getparent().remove(grupo._element)
        base_part = base.part
        sw, sh = int(prs.slide_width), int(prs.slide_height)
        # Rejilla: cuantas caben desde (x0,y0) con paso = tamaño de la etiqueta.
        cols = max(1, (sw - x0) // lw)
        rows = max(1, (sh - y0) // lh)
        posiciones = [(x0 + c * lw, y0 + r * lh)
                      for r in range(rows) for c in range(cols)]
        n_hojas = max(1, math.ceil(len(valores) / len(posiciones)))
        hojas = [base]
        for _h in range(1, n_hojas):
            hojas.append(self._etiqueta_clonar_slide(prs, base))
        self._etiqueta_tile(prs, base, base_part, unidad, posiciones,
                            hojas, 0, valores)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _etiqueta_construir_combinado_pptx(self, subtipo, datos, bloques,
                                           buffer_specs):
        """UN pptx con TODAS las etiquetas JUNTAS en las mismas hojas: primero
        las de CAJA (etiqueta grande, rejilla de 3 columnas) y, en el ESPACIO
        LIBRE que queda debajo, las de BUFFER (rejilla de 6 columnas). Si no
        caben, agrega las hojas tabloide necesarias, aprovechando el espacio.

        buffer_specs = lista de (plantilla, valores). Reusa el deck de caja
        (que ya deja libre la parte de abajo de la ultima hoja) y le agrega los
        buffers fluyendo hacia abajo desde la ultima etiqueta de caja."""
        import copy
        import io
        import os
        from pptx import Presentation
        from pptx.util import Inches

        # 1) Deck de caja (reusa el builder existente).
        caja_bytes = self._etiqueta_construir_pptx(subtipo, datos, bloques)
        prs = Presentation(io.BytesIO(caja_bytes))
        sheet_h = int(prs.slide_height)
        sheet_w = int(prs.slide_width)
        x0 = int(Inches(0.19))
        y0 = int(Inches(0.25))
        gap = int(Inches(0.08))
        tdir = os.path.join(os.path.dirname(__file__), '..', 'static', 'templates')

        def _hoja_en_blanco():
            # Clona la ultima hoja (conserva marco/fondo) y le quita las
            # etiquetas, para una hoja nueva solo con el marco tabloide.
            base = prs.slides[-1]
            nueva = self._etiqueta_clonar_slide(prs, base)
            for sh in list(nueva.shapes):
                if sh.shape_type == 6:
                    sh._element.getparent().remove(sh._element)
            return nueva

        # Punto de inicio de los buffers = justo debajo de la ultima etiqueta
        # de caja de la ultima hoja.
        cur = prs.slides[-1]
        grupos = [sh for sh in cur.shapes if sh.shape_type == 6]
        cur_y = (max(int(g.top) + int(g.height) for g in grupos) + gap) if grupos else y0

        for plantilla, valores in buffer_specs:
            if not valores:
                continue
            ruta = os.path.join(tdir, 'PLANTILLA_BUFFER_%s.pptx' % plantilla)
            if not os.path.exists(ruta):
                raise UserError(_(
                    'No existe la plantilla PLANTILLA_BUFFER_%s.pptx.') % plantilla)
            tmpl = Presentation(ruta)
            bslide = tmpl.slides[0]
            bgrupo = next((sh for sh in bslide.shapes if sh.shape_type == 6), None)
            if bgrupo is None:
                raise UserError(_(
                    'PLANTILLA_BUFFER_%s.pptx no tiene la etiqueta como grupo.')
                    % plantilla)
            lw, lh = int(bgrupo.width), int(bgrupo.height)
            unidad = copy.deepcopy(bgrupo._element)
            media = self._etiqueta_importar_medias(prs, bslide.part, bgrupo._element)
            cols = max(1, (sheet_w - x0) // lw)
            col_x = [x0 + c * lw for c in range(cols)]
            i = 0
            n = len(valores)
            while i < n:
                if cur_y + lh > sheet_h:
                    cur = _hoja_en_blanco()
                    cur_y = y0
                for c in range(cols):
                    if i >= n:
                        break
                    nuevo = copy.deepcopy(unidad)
                    self._etiqueta_mover_grupo(nuevo, col_x[c], cur_y)
                    self._etiqueta_llenar_element(nuevo, valores[i])
                    cur.shapes._spTree.append(nuevo)
                    if media:
                        self._etiqueta_remap_medias(nuevo, media, cur.part)
                    i += 1
                cur_y += lh

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()

    def _etiqueta_render_preview(self, pptx_bytes):
        """Renderiza un PPTX de UNA etiqueta a PNG (LibreOffice headless) y lo
        recorta al area de la etiqueta. Devuelve bytes PNG (o b'' si falla).
        Se usa para la vista previa de una sola etiqueta en el dialogo."""
        import io
        import os
        import shutil
        import subprocess
        import tempfile
        try:
            from PIL import Image
            from pptx import Presentation
        except Exception:
            return b''
        if not pptx_bytes:
            return b''
        # Bbox de la etiqueta = el (primer) grupo de la hoja.
        prs = Presentation(io.BytesIO(pptx_bytes))
        s = prs.slides[0]
        grupo = next((sh for sh in s.shapes if sh.shape_type == 6), None)
        if grupo is None:
            return b''
        sw, sh_ = float(prs.slide_width), float(prs.slide_height)
        gl, gt = float(grupo.left), float(grupo.top)
        gw, gh = float(grupo.width), float(grupo.height)
        # Pre-proceso SOLO para el render: a los cuadros cuyo texto es un unico
        # token (LOT, REF, caducidad, etiquetas cortas) se les quita el ajuste
        # de linea. Asi LibreOffice no parte "0726/01/TIF" en la diagonal. Los
        # cuadros con varias palabras (lista "Contiene") NO se tocan.
        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN

        def _no_wrap(shapes):
            for sh in shapes:
                if sh.shape_type == 6:
                    _no_wrap(sh.shapes)
                    continue
                if sh.has_text_frame:
                    t = (sh.text_frame.text or '').strip()
                    if t and not any(c.isspace() for c in t):
                        try:
                            sh.text_frame.auto_size = MSO_AUTO_SIZE.NONE
                        except Exception:
                            pass
                        try:
                            if t.upper() in ('LOT', 'REF'):
                                # Labels con recuadro. El cuadro es muy angosto y
                                # LibreOffice partiria "LOT"; se ENSANCHA el
                                # cuadro (centrado en su mismo centro) para que
                                # quepa en una linea y, con parrafo centrado,
                                # quede centrado sobre el recuadro. Wrap activo.
                                cx = int(sh.left) + int(sh.width) // 2
                                neww = int(sh.width) * 3
                                sh.width = neww
                                sh.left = cx - neww // 2
                                sh.text_frame.word_wrap = True
                                sh.text_frame.margin_left = 0
                                sh.text_frame.margin_right = 0
                                sh.text_frame.margin_top = 0
                                sh.text_frame.margin_bottom = 0
                                sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                                # Forzar centrado horizontal del parrafo (por si
                                # el algn del template no sobrevive el proceso).
                                for _p in sh.text_frame.paragraphs:
                                    _p.alignment = PP_ALIGN.CENTER
                            else:
                                # Valores (lote/caducidad/ref): NO cortar el token
                                # y anclar al fondo para quedar al mismo nivel.
                                sh.text_frame.word_wrap = False
                                if 'ontien' not in t.lower():
                                    sh.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
                        except Exception:
                            pass
        _no_wrap(s.shapes)
        tmpd = tempfile.mkdtemp(prefix='lblprev_')
        try:
            pptx_path = os.path.join(tmpd, 'l.pptx')
            prs.save(pptx_path)
            prof = 'file://' + os.path.join(tmpd, 'prof')
            # Render a alta resolucion (200 DPI) para que la etiqueta salga
            # nitida, no pixelada. PixelWidth/Height segun el tamano de la hoja.
            dpi = 200
            pxw = int(sw / 914400.0 * dpi)
            pxh = int(sh_ / 914400.0 * dpi)
            flt = ('png:impress_png_Export:{"PixelWidth":{"type":"long",'
                   '"value":%d},"PixelHeight":{"type":"long","value":%d}}'
                   % (pxw, pxh))
            subprocess.run(
                ['soffice', '--headless', '-env:UserInstallation=' + prof,
                 '--convert-to', flt, '--outdir', tmpd, pptx_path],
                capture_output=True, timeout=120)
            png_path = os.path.join(tmpd, 'l.png')
            if not os.path.exists(png_path):
                return b''
            im = Image.open(png_path)
            W, H = im.size
            mx, my = 0.02 * gw, 0.02 * gh
            box = (max(0, int((gl - mx) / sw * W)),
                   max(0, int((gt - my) / sh_ * H)),
                   min(W, int((gl + gw + mx) / sw * W)),
                   min(H, int((gt + gh + my) / sh_ * H)))
            crop = im.crop(box)
            # Letterbox a proporcion FIJA (3:2): se ESCALA la etiqueta (hacia
            # arriba o abajo) para llenar el lienzo respetando su proporcion, y
            # se centra en blanco. Asi el widget la muestra grande y sin deformar.
            CW, CH = 900, 600
            scale = min(CW / crop.width, CH / crop.height)
            new = crop.resize((max(1, int(crop.width * scale)),
                               max(1, int(crop.height * scale))), Image.LANCZOS)
            canvas = Image.new('RGB', (CW, CH), (255, 255, 255))
            ox, oy = (CW - new.width) // 2, (CH - new.height) // 2
            if new.mode == 'RGBA':
                canvas.paste(new, (ox, oy), new)
            else:
                canvas.paste(new.convert('RGB'), (ox, oy))
            buf = io.BytesIO()
            canvas.save(buf, format='PNG')
            return buf.getvalue()
        except Exception:
            return b''
        finally:
            shutil.rmtree(tmpd, ignore_errors=True)
